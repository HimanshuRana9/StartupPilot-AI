import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable
from app.models.schemas import UnifiedStartupReport

def generate_pdf_report(report: UnifiedStartupReport) -> bytes:
    """
    Generates an executive, publication-grade PDF startup evaluation report using ReportLab.
    """
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=letter,
        rightMargin=40,
        leftMargin=40,
        topMargin=40,
        bottomMargin=40
    )
    
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        'DocTitle',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=24,
        textColor=colors.HexColor('#4F46E5'),
        spaceAfter=6
    )
    subtitle_style = ParagraphStyle(
        'DocSubTitle',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=12,
        textColor=colors.HexColor('#6B7280'),
        spaceAfter=15
    )
    heading_style = ParagraphStyle(
        'SectionHeading',
        parent=styles['Heading2'],
        fontName='Helvetica-Bold',
        fontSize=14,
        textColor=colors.HexColor('#1E1B4B'),
        spaceBefore=12,
        spaceAfter=8
    )
    body_style = ParagraphStyle(
        'BodyTextCustom',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=10,
        textColor=colors.HexColor('#1F2937'),
        leading=14,
        spaceAfter=6
    )
    
    story = []
    
    # Title Header
    story.append(Paragraph("StartupPilot AI — Business Feasibility & Pitch Report", title_style))
    story.append(Paragraph(f"Startup Idea: <b>{report.idea}</b> | Overall Score: <b>{report.overall_readiness_score}/10</b> | Date: {report.timestamp}", subtitle_style))
    story.append(HRFlowable(width="100%", thickness=1.5, color=colors.HexColor('#6366F1'), spaceAfter=12))
    
    # Section 1: Executive Summary & Idea Validation
    story.append(Paragraph("1. Idea Validation & Executive Summary", heading_style))
    story.append(Paragraph(f"<b>Problem Statement:</b> {report.idea_validation.problem_statement}", body_style))
    story.append(Paragraph(f"<b>Proposed Solution:</b> {report.idea_validation.proposed_solution}", body_style))
    story.append(Paragraph(f"<b>Target Audience:</b> {report.idea_validation.target_audience}", body_style))
    story.append(Paragraph(f"<b>Initial Recommendation:</b> {report.idea_validation.initial_recommendation}", body_style))
    story.append(Spacer(1, 10))
    
    # Section 2: Market & Competitors
    story.append(Paragraph("2. Market Analysis & Competitor Landscape", heading_style))
    story.append(Paragraph(f"<b>Market Growth & Size:</b> {report.market_research.market_size_description} ({report.market_research.cagr_growth_rate})", body_style))
    story.append(Paragraph(f"<b>Market Gap:</b> {report.competitor_analysis.market_gap}", body_style))
    
    comp_table_data = [["Competitor Name", "Pricing Model", "Key Vulnerabilities"]]
    for comp in report.competitor_analysis.competitors:
        comp_table_data.append([comp.name, comp.pricing_model, ", ".join(comp.vulnerabilities)])
        
    t = Table(comp_table_data, colWidths=[150, 120, 260])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#EEF2FF')),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.HexColor('#3730A3')),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, -1), 9),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
        ('TOPPADDING', (0, 0), (-1, -1), 6),
        ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#E0E7FF')),
    ]))
    story.append(t)
    story.append(Spacer(1, 10))
    
    # Section 3: Human-Level Regional Supply Chain & Price Arbitrage
    story.append(Paragraph("3. Regional Supply Chain & Cost Arbitrage", heading_style))
    story.append(Paragraph(f"<b>Recommended Manufacturing/Operations Hub:</b> {report.regional_arbitrage.recommended_setup_location}", body_style))
    story.append(Paragraph(f"<b>Recommended Target Sales Hub:</b> {report.regional_arbitrage.recommended_sales_location}", body_style))
    
    for opp in report.regional_arbitrage.arbitrage_opportunities:
        story.append(Paragraph(f"• <b>{opp.source_location} ➔ {opp.target_location}</b>: {opp.strategic_advice} (Cost Delta: {opp.cost_difference_percent}, Profit Boost: {opp.estimated_profit_margin_boost})", body_style))
    story.append(Spacer(1, 10))
    
    # Section 4: Cost Estimation & Funding Roadmap
    story.append(Paragraph("4. Cost Estimation & Funding Strategy", heading_style))
    story.append(Paragraph(f"<b>Initial Capital Required:</b> ₹{report.cost_estimation.total_initial_budget_required:,.2f} | <b>Monthly Burn:</b> ₹{report.cost_estimation.monthly_burn_rate:,.2f}", body_style))
    story.append(Paragraph(f"<b>Primary Funding Advice:</b> {report.funding_advisor.primary_recommendation}", body_style))
    story.append(Paragraph(f"<b>Unit Economics:</b> {report.cost_estimation.unit_economics_summary}", body_style))
    story.append(Spacer(1, 10))
    
    # Section 5: Month-by-Month Execution Roadmap
    story.append(Paragraph("5. Implementation Roadmap", heading_style))
    for milestone in report.roadmap.milestones:
        tasks_str = "; ".join(milestone.key_tasks)
        story.append(Paragraph(f"<b>Month {milestone.month} ({milestone.phase_name}):</b> {tasks_str} [Output: {milestone.expected_output}]", body_style))
        
    doc.build(story)
    buffer.seek(0)
    return buffer.getvalue()

def generate_pptx_pitch_deck(report: UnifiedStartupReport) -> bytes:
    """
    Generates an 8-slide PowerPoint Pitch Deck (.pptx) using python-pptx with vibrant presentation styling.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # Helper to format title & content
    def add_styled_slide(title_text: str, bullets: list, takeaway: str = ""):
        slide = prs.slides.add_slide(blank_layout)
        
        # Header Background banner
        shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(13.333), Inches(1.2)  # MSO_SHAPE.RECTANGLE = 1
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(30, 27, 75)  # Dark Indigo
        shape.line.fill.background()
        
        # Header Title
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Arial"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Content Box
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.5))
        ctf = content_box.text_frame
        ctf.word_wrap = True
        
        for idx, bullet in enumerate(bullets):
            p = ctf.add_paragraph() if idx > 0 else ctf.paragraphs[0]
            p.text = f"•  {bullet}"
            p.font.name = "Arial"
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(31, 41, 55)
            p.space_after = Pt(14)
            
        # Takeaway Footer Box
        if takeaway:
            foot_shape = slide.shapes.add_shape(
                1, Inches(0.8), Inches(6.3), Inches(11.733), Inches(0.8)
            )
            foot_shape.fill.solid()
            foot_shape.fill.fore_color.rgb = RGBColor(238, 242, 255)
            foot_shape.line.color.rgb = RGBColor(99, 102, 241)
            
            tf_foot = foot_shape.text_frame
            p_foot = tf_foot.paragraphs[0]
            p_foot.text = f"Key Takeaway: {takeaway}"
            p_foot.font.name = "Arial"
            p_foot.font.size = Pt(14)
            p_foot.font.bold = True
            p_foot.font.color.rgb = RGBColor(67, 56, 202)
            p_foot.alignment = PP_ALIGN.LEFT
            
        return slide
        
    # Slide 1: Cover Slide
    cover_slide = prs.slides.add_slide(blank_layout)
    cover_bg = cover_slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    cover_bg.fill.solid()
    cover_bg.fill.fore_color.rgb = RGBColor(15, 23, 42)  # Dark slate
    
    tbox = cover_slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.333), Inches(3))
    tf = tbox.text_frame
    p1 = tf.paragraphs[0]
    p1.text = report.idea.title()
    p1.font.name = "Arial"
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(129, 140, 248)
    
    p2 = tf.add_paragraph()
    p2.text = "Investor Pitch Deck & Strategic Business Plan"
    p2.font.name = "Arial"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(226, 232, 240)
    p2.space_before = Pt(10)
    
    p3 = tf.add_paragraph()
    p3.text = f"Generated by StartupPilot AI | Overall Score: {report.overall_readiness_score}/10"
    p3.font.name = "Arial"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(148, 163, 184)
    p3.space_before = Pt(25)

    # Deck Slides from report
    for slide_data in report.investor_readiness.deck_slides:
        add_styled_slide(
            title_text=slide_data.title,
            bullets=slide_data.bullet_points,
            takeaway=slide_data.key_takeaway
        )
        
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()
